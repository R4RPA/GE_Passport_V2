import json

import utilities.support_functions as sup_functions
import collections
import collections.abc
from pptx import Presentation
import os.path


def create_ouput_pptx(data):
    print("========> create_ouput_pptx")
    prs = Presentation(data['ppt_template'])

    """Update title slide"""
    slide = prs.slides[0]
    sup_functions.replace_text_in_slide(slide, data)

    """Update summary slides 1 to 4"""
    slide = prs.slides[1]
    sup_functions.replace_text_in_slide(slide, data)
    data_row1 = 3
    data_row2 = 2
    data_row3 = 2
    line_num = 1
    for ha_file in data['ha_files']:
        ha_file['nc_data']['t1_ln'] = str(line_num).zfill(2)
        sup_functions.update_text_in_slide_tables(prs.slides[2], ha_file['nc_data'], 2, data_row1)
        data_row1 += 1
        if 'eval_data' in ha_file and len(ha_file['eval_data']) > 0:
            ha_file['eval_data']['t2_ln'] = str(line_num).zfill(2)
            sup_functions.update_text_in_slide_tables(prs.slides[3], ha_file['eval_data'], 1, data_row2)
            data_row2 += 1
        if 'eval_nc_data' in ha_file and len(ha_file['eval_nc_data']) > 0:
            ha_file['eval_nc_data']['t3_ln'] = str(line_num).zfill(2)
            sup_functions.update_text_in_slide_tables(prs.slides[4], ha_file['eval_nc_data'], 1, data_row3)
            data_row3 += 1
        line_num += 1

    """update nc image slide, set max 3 images per slide, insert number of slides required"""
    num_slides, nc_images_list = sup_functions.get_number_of_nc_images_slides_to_insert(data)
    prs = sup_functions.insert_nc_images_slides(prs, data, num_slides, nc_images_list)

    """update nc slides based on type identified """
    insert_nc_type_slides(prs, data, num_slides)

    prs.save(data['output_file'])


def insert_nc_type_slides(prs, data, num_slides):
    """map slide index range per nc type"""
    type_1_start_slide = 6 + (num_slides-1)
    type_1_end_slide = type_1_start_slide + 1  # 7
    type_2_start_slide = type_1_end_slide + 1  # 8
    type_2_end_slide = type_2_start_slide + 6  # 14
    type_3_start_slide = type_2_end_slide + 1  # 15
    type_3_end_slide = type_3_start_slide + 3  # 18
    dest_slide_index = type_3_end_slide + 1  # 19

    for ha_index, ha_file in enumerate(data['ha_files']):
        slide3 = None
        slide4 = None
        slide5 = None
        slide6 = None
        slide7 = None

        """select respective nc slides """
        if ha_file['template_type'] == '1':
            start = type_1_start_slide
            end = type_1_end_slide + 1
        elif ha_file['template_type'] == '2':
            start = type_2_start_slide
            end = type_2_end_slide + 1
        elif ha_file['template_type'] == '3':
            start = type_3_start_slide
            end = type_3_end_slide + 1
        start_index = dest_slide_index

        head_layout = prs.slide_masters[0].slide_layouts[1]
        data_layout = prs.slide_masters[0].slide_layouts[4]

        """Insert required nc slides"""
        for copy_slide_index in range(start, end):
            if copy_slide_index == start:
                block_layout = head_layout
            else:
                block_layout = data_layout
            sup_functions.duplicate_slide(prs, block_layout, copy_slide_index, dest_slide_index)
            dest_slide_index += 1

        """Update nc slides"""
        slide1 = prs.slides[start_index]
        slide2 = prs.slides[start_index+1]
        image_slides = []
        if 'nc_images' in ha_file:
            image_data = {'slide': slide2, 'images': ha_file['nc_images']}
            image_slides.append(image_data)

        """create a dictionary slide vs images to insert"""
        if ha_file['template_type'] in ['2', '3']:
            slide3 = prs.slides[start_index + 2]
            slide4 = prs.slides[start_index + 3]
            if ha_file['template_type'] == '2':

                image_data = {'slide': slide4, 'images': ha_file['mcb_data']}
                image_slides.append(image_data)

                slide5 = prs.slides[start_index + 4]
                image_data = {'slide': slide5, 'images': ha_file['fbo_data']}
                image_slides.append(image_data)

                slide6 = prs.slides[start_index + 5]
                image_data = {'slide': slide6, 'images': ha_file['fbo_data']}
                image_slides.append(image_data)

                slide7 = prs.slides[start_index + 6]

            if ha_file['template_type'] == '3':

                image_data = {'slide': slide4, 'images': ha_file['inspect_id_data']}
                image_slides.append(image_data)

        """update images in slides"""
        update_image_slide(image_slides)
        delete_palceholders_from_image_slide(image_slides)

        """update text in slides"""
        slides_list = [slide1, slide2, slide3, slide4, slide5, slide6, slide7]
        for slide in slides_list:
            if slide:
                sup_functions.replace_text_in_slide(slide, ha_file['nc_data'])
                if 'fbo_data' in ha_file:
                    sup_functions.replace_text_in_slide(slide, ha_file['fbo_data'])
                if 'mcb_data' in ha_file:
                    sup_functions.replace_text_in_slide(slide, ha_file['mcb_data'])
                sup_functions.replace_text_in_slide(slide, ha_file)
                sup_functions.replace_text_in_slide(slide, data)

                sup_functions.replace_text_in_slide_tables(slide, ha_file['nc_data'], font_size=8)

    """cleanup"""
    for index in range(type_1_start_slide, type_3_end_slide+1):
        sup_functions.delete_slide(prs, [type_1_start_slide])

    sup_functions.delete_table_rows(prs)

    return prs


def update_image_slide(image_slides):
    """update images in slides"""
    for image_data in image_slides:
        slide = image_data['slide']
        for key, image_path in image_data['images'].items():
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if 'image' in shape.text and key in shape.text:
                        slide.shapes.add_picture(image_path, shape.left, shape.top, shape.width, shape.height)
                        slide.shapes.element.remove(shape.element)


def delete_palceholders_from_image_slide(image_slides):
    """cleanup"""
    for image_data in image_slides:
        slide = image_data['slide']
        for shape in slide.shapes:
            if shape.has_text_frame:
                if '{image' in shape.text:
                    slide.shapes.element.remove(shape.element)


def main():
    input_path = r"C:\Users\HP\PycharmProjects\GE_Passport\input"
    output_file = os.path.join(input_path, "ref_data.json")
    with open(output_file, 'r') as f:
        data = json.load(f)

    create_ouput_pptx(data)



if __name__ == '__main__':

    main()
