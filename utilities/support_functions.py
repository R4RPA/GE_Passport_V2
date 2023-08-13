import math
import os, shutil
import glob
import copy
from datetime import datetime
from copy import deepcopy
import collections
import collections.abc
from pptx.util import Pt
from pptx.oxml.xmlchemy import OxmlElement
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.table import Table, _Row, _Column, _Cell
from pptx.dml.color import ColorFormat, RGBColor
import win32com.client as win32
from time import sleep
import ctypes
from PIL import ImageGrab
import re
import tempfile


def get_bullet_info(paragraph, run=None):
    """Get info about bullet text format to reapply to another part of text"""
    pPr = paragraph._p.get_or_add_pPr()
    if run is None:
        run = paragraph.runs[0]
    p_info = {"marL": pPr.attrib['marL'], "indent": pPr.attrib['indent'], "level": paragraph.level,
              "fontName": run.font.name, "fontSize": run.font.size, }
    return p_info


def _sub_element(parent, tagname, **kwargs):
    """Helper for Paragraph bullet Point"""
    element = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element


def set_bullet_paragraph(paragraph, bullets_info):
    """Apply bullets text format for given paragraph"""
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set('marL', bullets_info['marL'])
    pPr.set('indent', bullets_info['indent'])
    _ = _sub_element(parent=pPr, tagname="a:buSzPct", val="100000")
    _ = _sub_element(parent=pPr, tagname="a:buFont", typeface=bullets_info['fontName'])
    _ = _sub_element(parent=pPr, tagname='a:buChar', char="•")


def move_slide(prs, old_index, new_index):
    """Move or re-arrange the slide to desired position"""
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])


def check_slide_layouts(prs):
    dest_slide_index = 0
    slide_num = 0
    for num in range(100):
        for num2 in range(100):
            try:
                print('TRY slide_num', slide_num, 'num', num, 'num2', num2)
                nc_layout = prs.slide_masters[num].slide_layouts[num2]
                copy_slide_index = 0
                dest_slide_index += 1
                duplicate_slide(prs, nc_layout, copy_slide_index, dest_slide_index)
                slide_num += 1

                print('DONE slide_num', slide_num, 'num', num, 'num2', num2, 'dest_slide_index', dest_slide_index)
            except Exception as e:
                print(str(e))
                pass


def duplicate_slide(prs, layout, copy_slide_index, dest_slide_index):
    """Duplicate the slide with the given index in presentation.
    Adds slide to the index + 3 position of the presentation"""

    """Original slide"""
    source = prs.slides[copy_slide_index]
    """Add a blank slide with selected layout"""
    dest = prs.slides.add_slide(layout)
    slide_id = prs.slides.index(dest)
    """Move or re-arrange the slide to desired position"""
    move_slide(prs, slide_id, dest_slide_index)

    """delete any default shapes are in the slide"""
    for shape in dest.shapes:
        sp = shape._sp
        sp.getparent().remove(sp)

    """Copy each shape in the original slide to duplicated slide"""
    for shape in source.shapes:

        if 'Picture' in shape.name:
            """if shape has image, save the image to local, add to slide in position and delete"""
            dummy_image = shape.name + '.jpg'
            with open(dummy_image, 'wb') as f:
                f.write(shape.image.blob)
            dest.shapes.add_picture(dummy_image, shape.left, shape.top, shape.width, shape.height)
            os.remove(dummy_image)
        else:
            newel = copy.deepcopy(shape.element)
            dest.shapes._spTree.insert_element_before(newel, 'p:extLst')
            
    #return prs


def delete_slide(prs, delete_slide_index: list):
    """Delete a slide from presentation"""
    xml_slides = prs.slides._sldIdLst  
    slides = list(xml_slides)
    for index in delete_slide_index:
        xml_slides.remove(slides[index])
    
    return prs


def get_object_by_name(slide, object_name):
    """Search and return the shape by name"""
    for shape in slide.shapes:
        if object_name in shape.name:
            if 'Table' in object_name:
                return shape.table
            else:
                return shape
    return None


def add_multiple_rows(table: Table, number_of_rows: int = 1, add_empty_row: bool = False):
    """Add multiple rows into table object"""
    for row in range(number_of_rows):
        table = add_row(table, add_empty_row)
    return table


def add_row(table: Table, add_empty_row: bool = False):
    """Add single row into table object"""
    new_row = deepcopy(table._tbl.tr_lst[-1])
    if add_empty_row:
        """if add_empty_row == True, then clear the cell content"""
        for tc in new_row.tc_lst:
            cell = _Cell(tc, new_row.tc_lst)
            cell.text = ''
    table._tbl.append(new_row)
    return table

def delete_table_rows(prs):
    """Delete table rows"""
    for slide in prs.slides:
        for shape in slide.shapes:
            try:
                table = shape.table
                row_count = len(table.rows)
                for row in range(row_count, 0, -1):
                    try:
                        cell = table.cell(row, 1)
                        cell_text = cell.text
                        if cell_text.startswith('{') or cell_text == '':
                            remove_row(table, table.rows[row])
                    except:
                        pass
            except:
                pass

def remove_row(table: Table, row_to_delete: _Row):
    """Delete Table Row"""
    table._tbl.remove(row_to_delete._tr)


def replace_text_in_slide_tables(slide, nc_result: dict, font_size=12):
    """Replace tags in slide with nc results"""
    for shape in slide.shapes:
        try:
            table = shape.table
            update_table_values(table, nc_result, font_size)
        except:
            pass


def update_table_values(table: Table, nc_result: dict, font_size=12):
    row_count = len(table.rows)
    for row in range(-1, row_count-1):
        update_table_row(table, nc_result, row, font_size)


def update_table_row(table: Table, nc_result: dict, row_num: int = 0, font_size=12):
    """Update table row with nc results by tag identified"""
    for column in range(len(table.columns)):
        cell = table.cell(row_num + 1, column)
        cell_text = cell.text
        for match, replacement in nc_result.items():
            match = '{' + match + '}'
            if match in cell_text and replacement != '':
                cell_text = cell_text.replace(str(match), str(replacement))
                cell.text = cell_text
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(font_size)
                cell.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                cell.text_frame.word_wrap = False
    return table


def update_text_in_slide_tables(slide, nc_result: dict, head_row, data_row):
    """Replace tags in slide with nc results"""
    for shape in slide.shapes:
        try:
            table = shape.table
            update_table_row_v2(table, nc_result, head_row, data_row)
        except:
            pass
    return slide

def update_table_row_v2(table: Table, nc_result: dict, head_row:int, data_row: int, font_size=8):
    """Update table row with nc results by tag identified"""
    for column in range(len(table.columns)):
        header_cell = table.cell(head_row, column)
        data_cell = table.cell(data_row, column)
        for match, replacement in nc_result.items():
            match = '{' + match + '}'
            if match in header_cell.text and replacement != '':
                data_cell.text = replacement
                for paragraph in data_cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(font_size)
                data_cell.text_frame.word_wrap = True
                data_cell.text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    return table

def replace_text_in_prs(prs, nc_result: dict):
    for slide in prs.slides:
        replace_text_in_slide(slide, nc_result)


def replace_text_in_slide(slide, nc_result: dict):
    """Replace tags in slide with nc results"""
    green_txt = 'margin is in acceptable limits. Hence, this NC can be accepted from clearance standpoint.'
    red_txt = 'margin is not in acceptable limits. Hence, additional study is required.'
    for shape in slide.shapes:
        replace_text_in_shape(shape, nc_result)
        if shape.has_text_frame:
            if shape.text and green_txt in shape.text:
                shape.fill.solid()
                shape.fill.fore_color.rgb  = RGBColor(200, 230, 170)
            elif shape.text and red_txt in shape.text:
                shape.fill.solid()
                shape.fill.fore_color.rgb  = RGBColor(250, 200, 200)
    return slide
                
def replace_text_in_shape(shape, nc_result: dict):
    """Replace tag text in a shape"""
    for match, replacement in nc_result.items():
        match = '{' + match + '}'
        if shape.has_text_frame:
            if (shape.text.find(match)) != -1:
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    whole_text = "".join(run.text for run in paragraph.runs)
                    whole_text = whole_text.replace(str(match), str(replacement))
                    for idx, run in enumerate(paragraph.runs):
                        if idx != 0:
                            p = paragraph._p
                            p.remove(run._r)
                    if bool(paragraph.runs):
                        font_size = paragraph.runs[0].font.size
                        paragraph.runs[0].text = whole_text
                        paragraph.runs[0].font.color.rgb = RGBColor(0x0, 0x0, 0x0)
                        paragraph.runs[0].font.size = font_size

def replace_images_in_slide(slide, nc_result: dict):
    
    """Replace images in slide"""
    for shape in slide.shapes:
        for key in nc_result:
            if 'Image' in key or  'Score' in key:
                image_name = key.replace('{','').replace('}','')
                image_path = nc_result[key]
                if len(image_path) > 4 and image_name in shape.name:
                    slide.shapes.add_picture(image_path, shape.left, shape.top, shape.width, shape.height)
                    slide.shapes.element.remove(shape.element)   


def delete_prev_nc_images(nc_result):
    for key in nc_result:
        if 'Image' in key or 'Score' in key:
            image_path = nc_result[key]
            try:
                os.remove(image_path)
            except Exception as e:
                pass

def get_number_of_nc_images_slides_to_insert_old(data):
    number_of_ha_files = len(data['ha_files'])
    return math.ceil(number_of_ha_files / 3)


def get_number_of_nc_images_slides_to_insert(data):
    number_of_ha_images = 0
    nc_images_list =[]
    for haf in data['ha_files']:
        if 'nc_images' in haf and len(haf['nc_images']) > 0:
            number_of_ha_images += 1
            nc_images_list.append(haf['nc_images']['1'])
    number_of_slides = math.ceil(number_of_ha_images / 3)
    return number_of_slides, nc_images_list

def insert_nc_images_slides(prs, data, num_slides, nc_images_list):
    copy_slide_index = 5
    dest_slide_index = copy_slide_index + 1
    block_layout = prs.slide_masters[0].slide_layouts[4]
    if num_slides > 1:
        for x in range(num_slides-1):
            duplicate_slide(prs, block_layout, copy_slide_index, dest_slide_index)
            dest_slide_index += 1

    for slide_idx in range(copy_slide_index, copy_slide_index + num_slides):
        image_list_for_slide = nc_images_list[(slide_idx - copy_slide_index) * 3: (slide_idx - copy_slide_index + 1) * 3]
        slide = prs.slides[slide_idx]
        text_list = ['{image1}', '{image2}', '{image3}']
        for index, image_path in enumerate(image_list_for_slide):
            image_text = text_list[index]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if shape.text == image_text:
                        slide.shapes.add_picture(image_path, shape.left, shape.top, shape.width, shape.height)

        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.text in text_list:
                    slide.shapes.element.remove(shape.element)

        replace_text_in_slide(slide, data)
    return prs


def get_column_number_from_range(cell_range, search_value):
    for cell in cell_range:
        cell_value = "".join(re.findall("[a-zA-Z0-9-+]+", str(cell.Value)))
        search_value = "".join(re.findall("[a-zA-Z0-9-+]+", str(search_value)))
        if cell_value.upper() == search_value.upper():
            return cell.Column
    return 0
    
def is_float(string):
    try:
        float(string)
        return True
    except ValueError:
        return False

def get_nc_input_docs(folder_path):
    """Search and list pdf files in the given folder"""
    nc_pdf_docs = glob.glob(os.path.join(folder_path, "*.pdf"))
    nc_pptx_docs = glob.glob(os.path.join(folder_path, "*.pptx"))
    nc_xlsx_docs = glob.glob(os.path.join(folder_path, "*.xlsx"))
    return nc_pdf_docs + nc_pptx_docs + nc_xlsx_docs


def check_for_files(folder_path, pattern1, pattern2=None, pattern3=None):
    """Search and list pdf files in the given folder"""
    list1 = glob.glob(os.path.join(folder_path, pattern1)) if pattern1 else []
    list2 = glob.glob(os.path.join(folder_path, pattern2)) if pattern2 else []
    list3 = glob.glob(os.path.join(folder_path, pattern3)) if pattern3 else []
    return list1 + list2 + list3


def get_output_nc_report_path(source_file_path, output_path):
    """Set Output file name"""
    today_date = datetime.now().strftime('%Y%m%d%H%M%S_')
    clean_report_name = os.path.basename(source_file_path).replace('.', "_").replace("  ", " ").replace(" ", "_")
    output_nc_filename = today_date + 'NC_Report_for_' + clean_report_name + '.pptx'
    output_file_path = os.path.join(output_path, output_nc_filename)
    return output_file_path


def killexcel():
    # Kill Excel if filed othwewise
    try:
        os.system('taskkill /IM EXCEL.exe /T /F')
    except:
        True
    sleep(5)
    
def get_excel():
    loop_count = 0
    while True:
        loop_count += 1
        try:
            """Initiate Excel Application"""
            killexcel()
            excel = win32.gencache.EnsureDispatch('Excel.Application')
            try:
                excel.Visible = True
                excel.AskToUpdateLinks = False
                excel.DisplayAlerts = False
                excel.CutCopyMode = False
                excel.Calculation = win32.constants.xlCalculationManual
            except:
                pass
            return excel
        except:
            delete_temp_files()
            killexcel()
            if loop_count > 5:
                raise ValueError('unable to start excel application')


class ExcelWrapper:
    def __enter__(self):
        self.excel = win32.gencache.EnsureDispatch('Excel.Application')
        try:
            self.excel.Visible = True
            self.excel.AskToUpdateLinks = False
            self.excel.DisplayAlerts = False
            self.excel.CutCopyMode = False
            self.excel.Calculation = win32.constants.xlCalculationManual
        except:
            pass
        return self.excel

    def __exit__(self, exc_type, exc_value, traceback):
        self.excel.Quit()

def delete_temp_files():
    try:
        temp_folder = os.environ['TEMP']
        for filename in os.listdir(temp_folder):
            file_path = os.path.join(temp_folder, filename)
            try:
                if os.path.isfile(file_path) or os.path.islink(file_path):
                    os.unlink(file_path)
                elif os.path.isdir(file_path):
                    shutil.rmtree(file_path)
            except:
                pass
    except:
        pass


def open_workbook(excel, xls_file):
    loop_count = 0
    while True:
        loop_count += 1
        try:
            wb = excel.Workbooks.Open(xls_file, None, False)

            if not wb:
                raise ValueError('Unable to open WB "{}", Retry again'.format(xls_file))
            return wb
        except Exception as e:
            print('open wb error ', str(e))
            killexcel()
            excel = get_excel()
            if loop_count > 5:
                raise ValueError('unable to open excel file {}'.format(xls_file))


def close_workbook(wb):
    try:
        wb.Close(False)
    except Exception as e:
        print(str(e))


def get_worksheet(wb, sheet_name):
    try:
        return wb.Sheets(sheet_name)
    except:
        return []

def get_worksheet_by_partial_text(wb, partial_str):
    for m, sheet in enumerate(wb.Sheets):
        if partial_str in sheet.Name:
            return sheet
    return None

def save_shape_as_image(shape, img_name):
    temp_dir = tempfile.gettempdir()
    file_path = os.path.join(temp_dir, img_name)
    loop_count = 0
    while True:
        loop_count += 1
        try:
            shape.Copy()
            image = ImageGrab.grabclipboard()
            if image.mode in ("RGBA", "P"):
                image = image.convert("RGB")
            image.save(file_path, 'png')
            return True, file_path
        except Exception as e:
            print("Unable to save object '{}' as image. Retry Loop {}".format(file_path, loop_count))
            clear_clip_board()
            sleep(1)
            if loop_count > 5:
                break
    return False, file_path


def save_range_as_image(copy_range, img_name):
    temp_dir = tempfile.gettempdir()
    file_path = os.path.join(temp_dir, img_name)
    loop_count = 0
    while True:
        loop_count += 1
        try:
            
            copy_range.Copy()
            image = ImageGrab.grabclipboard()
            if image.mode in ("RGBA", "P"):
                image = image.convert("RGB")
            image.save(file_path, 'png')
            return True, file_path
        except Exception as e:
            print("Unable to save object '{}' as image. Retry Loop {}".format(file_path, loop_count))
            clear_clip_board()
            sleep(1)
            if loop_count > 5:
                break
    return False, file_path


def isDigit(x):
    try:
        float(x)
        return True
    except:
        return False

def close_wb(wb):
    loop_count = 0
    while True:
        loop_count += 1
        try:
            wb.Close(False)
            break
        except:
            sleep(5)
            if loop_count > 5:
                break


def get_clearance_text(ClearanceTag):
    ClearanceDict = {'A': 'Axial',
                     'R': 'Radial',
                     'O': 'Overlap',
                     'C': 'Clashing'}
    if ClearanceTag in ClearanceDict:
        val = ClearanceDict[ClearanceTag]
    else:
        val = ClearanceTag
    return val


def quit_excel(excel):
    """Close Excel Application"""
    try:
        excel.Quit()
    except:
        killexcel()


def chck_rep1_nc_content(nc_dict):
    if len(nc_dict['NC_Slides']) > 0:
        for block in nc_dict['NC_Slides']:
            if len(nc_dict['NC_Slides'][block]['NC_Content']) > 0:
                return True
    return False
                        


def check_nc_dict_has_locations(nc_dict):
    try:
        nc_findings = ''
        for item in nc_dict["NC_Lines"]:
            Enum = item["{EngineNum}"]
            Snum = item["{StageNum}"]
            if item["{EffectedLocations}"]:
                return '', True
            nc_findings += "(For Engine '{}', Stage '{}' and No Effected Locations Found), ".format(Enum, Snum)
    except:
        pass
    return nc_findings, False


def clear_clip_board():
    loop_count = 0
    while True:
        loop_count += 1
        try:
            ctypes.windll.user32.OpenClipboard(None)
            ctypes.windll.user32.EmptyClipboard()
            ctypes.windll.user32.CloseClipboard()
            break
        except:
            sleep(5)
            if loop_count > 5:
                break


def get_part_number_by_jap_sting(wb):
    part_number = ''
    for m, sheet in enumerate(wb.Sheets):
        string = sheet.Cells(1, 1).Value
        if string and starts_with_japanese_char(string):
            part_number = string.split('=')[1]
            break

    return part_number

def starts_with_japanese_char(text):
    if len(text) > 0:
        first_char = text[0]
        # Check if the Unicode code point of the first character is within the Japanese ranges
        if ord('\u3040') <= ord(first_char) <= ord('\u309F') or \
           ord('\u30A0') <= ord(first_char) <= ord('\u30FF') or \
           ord('\u4E00') <= ord(first_char) <= ord('\u9FFF'):
            return True
    return False